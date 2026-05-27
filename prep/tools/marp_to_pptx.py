#!/usr/bin/env python3
"""Convert a Marp markdown slide deck to PowerPoint using a reference PPTX template.

Usage:
    python marp_to_pptx.py <input.md> <output.pptx> <template.pptx>
    python marp_to_pptx.py <input.md> <output.pptx>          # uses default template
    python marp_to_pptx.py                                    # uses all defaults

Template layout mapping (Master 4):
    0  - Title Slide           54 - Section Title
    4  - Title and Content     70 - Closing logo slide
   13  - Title Only            71 - Black Notes slide Layout
   45  - Developer Code Layout
"""

import argparse
import os
import re
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Default paths (relative to this script)
# ---------------------------------------------------------------------------
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.dirname(os.path.dirname(SCRIPT_DIR))
DEFAULT_TEMPLATE = os.path.join(SCRIPT_DIR, "template.pptx")
DEFAULT_INPUT = os.path.join(REPO_ROOT, "1-hour", "slides.md")
DEFAULT_OUTPUT = os.path.join(REPO_ROOT, "1-hour", "slides.pptx")

# ---------------------------------------------------------------------------
# Style constants
# ---------------------------------------------------------------------------
BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK = RGBColor(0x09, 0x1F, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

HEADING_FONT = "Segoe UI Semibold"
BODY_FONT = "Segoe UI"
CODE_FONT = "Consolas"

SLIDE_W = 12192000  # EMU
SLIDE_H = 6858000

# ---------------------------------------------------------------------------
# Layout mappings — match the branded corporate template used in the reference
# deck. The template has layouts spread across multiple masters.
# ---------------------------------------------------------------------------
# Master 0: Session opening layouts
#   Layout 0: "3.1 Session Divider" — opening title slide
#   Layout 1: "GREEN BLANK _ TITLE ONLY" — main content slides
# Master 1:
#   Layout 1: "7.1 Stat Slide" — section divider with big text/stats
# Master 2:
#   Layout 0: "3.1 Session Divider"
#   Layout 1: "GREEN BLANK _ TITLE ONLY"
#   Layout 2: "Divider and closing" — transition/conceptual slides
# Master 4:
#   Layout 71: "Black Notes slide Layout" — presenter notes

# We use the following layout strategy to match the reference deck:
MASTER_CONTENT = 0       # Master for GREEN BLANK _ TITLE ONLY
LAYOUT_CONTENT_IDX = 1   # "GREEN BLANK _ TITLE ONLY"
LAYOUT_TITLE_IDX = 0     # "3.1 Session Divider" (opening title)

MASTER_SECTION = 1       # Master for 7.1 Stat Slide
LAYOUT_SECTION_IDX = 1   # "7.1 Stat Slide" (section dividers)

MASTER_DIVIDER = 2       # Master for Divider and closing
LAYOUT_DIVIDER_IDX = 2   # "Divider and closing" (conceptual/transition)

MASTER_NOTES = 4         # Master for Black Notes slide Layout
LAYOUT_NOTES_IDX = 71    # "Black Notes slide Layout"
LAYOUT_CLOSING_IDX = 70  # "Closing logo slide"

# ---------------------------------------------------------------------------
# Slide-title → image file mapping
# ---------------------------------------------------------------------------
IMAGE_MAP = {
    "Azure Landing Zone — Conceptual Architecture": "alz-architecture.png",
    "Azure Landing Zone Architecture": "alz-architecture.png",
    "Azure Landing Zone Architecture — Management Group Hierarchy": "alz-architecture.png",
    "CAF for AI Adoption": "ai-adoption-process.png",
    "Landing Zone Design Areas": "alz-design-areas.png",
    "Landing Zone Design Areas for AI": "alz-design-areas.png",
    "Hub-Spoke Network Architecture for AI": "hub-spoke.png",
    "Hub-Spoke VNET Topology for AI": "hub-spoke.png",
    "Network Architecture — AI Landing Zone": "hub-spoke.png",
    "AI Application Landing Zone": "aoai-baseline-lz.png",
    "Platform vs Application Landing Zones": "alz-platform-app.png",
    "Platform Landing Zone vs Application Landing Zone": "alz-platform-app.png",
    "The Metropolis Analogy — Shared Platform Services": "metropolis.jpg",
    "Microsoft Defender for Cloud Overview": "security.jpg",
    "Microsoft Defender for Cloud — AI Protection": "security.jpg",
    "AI Operating Model": "cloud-infra.jpg",
    "AI Operating Model Overview": "cloud-infra.jpg",
    # Topic 1 - AI CoE & Governance concept diagrams
    "CoE Maturity Levels": "coe-maturity-levels.png",
    "Three CoE Operating Models": "coe-operating-models.png",
    "Agent Intake and Onboarding Process": "agent-intake-process.png",
    "Model Lifecycle Management": "model-lifecycle.png",
    "The Six Phases of AI Adoption": "six-phases-ai-adoption.png",
    "AI Governance Framework": "ai-governance-framework.png",
    "Community of Practice Model": "community-of-practice.png",
    "Onboarding a New Model": "onboarding-new-model.png",
    # Topic 1 - Visual infographics (replacing tables)
    "The POC Graveyard Problem": "poc-graveyard-failures.png",
    "Core CoE Roles": "coe-roles-cards.png",
    "Team Structure by Organisation Size": "coe-team-structure.png",
    "CoE Reporting Structure": "coe-reporting-structure.png",
    "RACI - AI Initiative Governance": "raci-governance.png",
    "Use-Case Intake and Prioritisation": "agent-intake-process.png",
    "Microsoft AI Technology Decision Tree": "ai-decision-tree.png",
    "POC-to-Production Graduation Checklist": "graduation-checklist.png",
    "Business Value Assessment Framework": "bva-framework.png",
    "Business Value Reporting": "bva-reporting.png",
    "FinOps for AI": "finops-cost-drivers.png",
    "AI Unit Economics": "unit-economics.png",
    "Role-Based Skilling Pathways": "skilling-pathways.png",
    "Microsoft Certification Roadmap": "certification-roadmap.png",
    "Continuous Learning Model": "continuous-learning.png",
    "Responsible AI - Six Principles Applied to Agents": "rai-six-principles.png",
    "OWASP Agentic AI Top 10 (2026)": "owasp-top10.png",
    "Agent Governance Toolkit - Runtime Security for Agents": "governance-toolkit.png",
    "Agent Governance - Implementation Roadmap": "implementation-roadmap.png",
    "Agent Governance - Responsibility Matrix": "governance-responsibility.png",
    "Agent 365 - Enterprise Control Plane": "agent365-capabilities.png",
    "Regulatory Compliance Timeline - 2026": "compliance-timeline.png",
    "AI Portfolio Management": "portfolio-metrics.png",
    "CoE Maturity Self-Assessment": "maturity-self-assessment.png",
    "Platform Selection Guide": "platform-guide.png",
    "Intelligence Integrations - Work IQ, Fabric IQ, Foundry IQ": "iq-integrations.png",
    "What to Do Next": "action-plan.png",
    # Topic 3 - Agent Architectures diagrams
    "Three-Plane Architecture Overview": "hub-centralization-vs-federation-matrix.png",
    "APIM AI Gateway - Token Management Detail": "apim-token-management-before-after.png",
    "Microsoft Fabric - Intelligent Data Platform": "intelligent-data-platform.png",
    "Agent Identity in Microsoft Entra ID": "agent-identity-entra-flow.png",
    "Networking Posture in the Landing Zone": "landing-zone-networking-posture.png",
    # Topic 4 - AgentOps diagrams
    "The AgentOps Lifecycle": "agentops-lifecycle-cycle.png",
    "Three Stages of Evaluation": "evaluation-three-stages-timeline.png",
    "Evaluation in CI/CD - Quality Gates": "ci-cd-evaluation-gate-pipeline.png",
    "The Observability Stack": "observability-stack-three-pillars.png",
    "Model Lifecycle Management": "model-lifecycle-upgrade-process.png",
    # Topic 5 - AI Security diagrams
    "Security Controls - Full Stack View": "security-full-stack-architecture.png",
    "Identity Hierarchy for Agents": "identity-hierarchy-agent-types.png",
    "Row-Level Security and Grounding Controls": "row-level-security-data-access.png",
    # Topic 6 - Solution Optimisation diagrams
    "Caching Strategies": "caching-strategy-three-levels.png",
    "Cost Modelling for Agentic Workloads": "cost-modeling-formula-breakdown.png",
    "Multi-Region Architecture for AI": "multi-region-architecture-diagram.png",
    "A/B Testing Agent Configurations": "model-upgrade-phased-rollout.png",
}

# Which images are architecture diagrams (replace code block, centered)
DIAGRAM_IMAGES = {
    "alz-architecture.png", "ai-adoption-process.png", "alz-design-areas.png",
    "hub-spoke.png", "aoai-baseline-lz.png", "alz-platform-app.png",
    "coe-maturity-levels.png", "coe-operating-models.png", "agent-intake-process.png",
    "model-lifecycle.png", "six-phases-ai-adoption.png", "ai-governance-framework.png",
    "community-of-practice.png", "onboarding-new-model.png",
    # Topic 1 - Visual infographics
    "poc-graveyard-failures.png", "coe-roles-cards.png", "coe-team-structure.png",
    "coe-reporting-structure.png", "raci-governance.png", "ai-decision-tree.png",
    "graduation-checklist.png", "bva-framework.png", "bva-reporting.png",
    "finops-cost-drivers.png", "unit-economics.png", "skilling-pathways.png",
    "certification-roadmap.png", "continuous-learning.png", "rai-six-principles.png",
    "owasp-top10.png", "governance-toolkit.png", "implementation-roadmap.png",
    "governance-responsibility.png", "agent365-capabilities.png", "compliance-timeline.png",
    "portfolio-metrics.png", "maturity-self-assessment.png", "platform-guide.png",
    "iq-integrations.png", "action-plan.png",
    # Topic 3
    "hub-centralization-vs-federation-matrix.png", "apim-token-management-before-after.png",
    "intelligent-data-platform.png", "agent-identity-entra-flow.png",
    "landing-zone-networking-posture.png", "interactive-assistant-token-flow.png",
    # Topic 4
    "agentops-lifecycle-cycle.png", "evaluation-three-stages-timeline.png",
    "ci-cd-evaluation-gate-pipeline.png", "observability-stack-three-pillars.png",
    "model-lifecycle-upgrade-process.png",
    # Topic 5
    "security-full-stack-architecture.png", "identity-hierarchy-agent-types.png",
    "row-level-security-data-access.png",
    # Topic 6
    "caching-strategy-three-levels.png", "cost-modeling-formula-breakdown.png",
    "multi-region-architecture-diagram.png", "model-upgrade-phased-rollout.png",
    # 1-hour deck custom visuals (code-rendered + AI)
    "complexity-ladder.png", "operating-loop.png", "maturity-ribbon.png",
    "foundry-control-plane.png", "telemetry-to-action.png",
    "red-teaming-taxonomy.png", "day2-quadrant.png",
    "production-gap.png", "start-with-one.png",
}


def _image_for_slide(title: str) -> str | None:
    """Return the image filename mapped to *title*, or None."""
    return IMAGE_MAP.get(title)


def _is_diagram(filename: str) -> bool:
    return filename in DIAGRAM_IMAGES


def _add_image_to_slide(slide, image_path: str, left, top, width=None, height=None,
                        alt_text: str = ""):
    """Add an image to the slide with accessibility alt text. If only width is given,
    height is computed from the image's aspect ratio (and vice-versa)."""
    if not os.path.isfile(image_path):
        return None
    if width and not height:
        with Image.open(image_path) as img:
            w, h = img.size
        aspect = h / w
        height = int(width * aspect)
    elif height and not width:
        with Image.open(image_path) as img:
            w, h = img.size
        aspect = w / h
        width = int(height * aspect)
    pic = slide.shapes.add_picture(image_path, left, top, width, height)
    # Set accessibility alt text (Microsoft accessibility guidelines)
    if alt_text:
        pic._element.nvPicPr.cNvPr.set("descr", alt_text)
        pic.name = alt_text
    return pic


# ---------------------------------------------------------------------------
# Markdown parsing
# ---------------------------------------------------------------------------
def _extract_speaker_notes(text: str) -> str:
    m = re.search(r"<!--\s*Speaker notes:\s*(.*?)\s*-->", text, re.DOTALL)
    return m.group(1).strip() if m else ""


def _is_lead(text: str) -> bool:
    return "<!-- _class: lead -->" in text


def _extract_table(lines: list) -> list | None:
    table_lines = []
    for ln in lines:
        stripped = ln.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            table_lines.append(stripped)
        elif table_lines:
            break
    if len(table_lines) < 2:
        return None
    rows = []
    for tl in table_lines:
        cells = [c.strip() for c in tl.strip("|").split("|")]
        if all(re.match(r"^[-:]+$", c) for c in cells):
            continue
        rows.append(cells)
    return rows if rows else None


def _extract_code_block(text: str) -> str | None:
    m = re.search(r"```[^\n]*\n(.*?)```", text, re.DOTALL)
    return m.group(1).rstrip() if m else None


def _extract_inline_images(text: str) -> list:
    """Extract markdown image references ![alt](path) from slide text."""
    return re.findall(r"!\[([^\]]*)\]\(([^)]+)\)", text)


def _extract_bullets(lines: list) -> list:
    bullets = []
    for ln in lines:
        stripped = ln.strip()
        if re.match(r"^[-*]\s", stripped):
            bullets.append(re.sub(r"^[-*]\s+", "", stripped))
    return bullets


def _extract_body_text(lines: list) -> list:
    body = []
    in_code = False
    in_table = False
    for ln in lines:
        stripped = ln.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            in_table = True
            continue
        if in_table and not stripped.startswith("|"):
            in_table = False
        if in_table:
            continue
        if stripped.startswith("#") or re.match(r"^[-*]\s", stripped):
            continue
        if stripped.startswith("<!--") or stripped == "":
            continue
        # Skip image references - they are handled separately
        if re.match(r"^!\[", stripped):
            continue
        body.append(stripped)
    return body


def parse_marp(md_text: str) -> list:
    """Parse Marp markdown into a list of slide dicts."""
    raw_blocks = re.split(r"\n---\s*\n", md_text)
    # Drop YAML front matter
    if raw_blocks:
        first = raw_blocks[0].strip()
        if first.startswith("---"):
            first = first.lstrip("-").strip()
        if first.startswith("marp:") or first.startswith("theme:"):
            raw_blocks = raw_blocks[1:]

    slides = []
    for idx, block in enumerate(raw_blocks):
        block = block.strip()
        if not block:
            continue

        speaker_notes = _extract_speaker_notes(block)
        is_lead = _is_lead(block)

        clean = re.sub(r"<!--.*?-->", "", block, flags=re.DOTALL).strip()
        if not clean:
            continue

        lines = clean.split("\n")
        title = ""
        subtitle = ""
        content_lines = []
        for ln in lines:
            stripped = ln.strip()
            if stripped.startswith("## "):
                subtitle = stripped.lstrip("#").strip()
            elif stripped.startswith("# "):
                title = stripped.lstrip("#").strip()
            else:
                content_lines.append(ln)

        table = _extract_table(content_lines)
        code_block = _extract_code_block(clean)
        bullets = _extract_bullets(content_lines)
        body_text = _extract_body_text(content_lines)
        inline_images = _extract_inline_images(clean)

        lower_title = title.lower().strip()
        if lower_title in ("thank you", "thanks", "thank you!"):
            stype = "thankyou"
        elif lower_title == "notes for presenter":
            stype = "notes"
        elif is_lead:
            stype = "title" if idx == 0 else "lead"
        elif idx == 0 and title and subtitle:
            stype = "title"
        else:
            stype = "content"

        slides.append(
            {
                "type": stype,
                "title": title,
                "subtitle": subtitle,
                "bullets": bullets,
                "table": table,
                "code_block": code_block,
                "speaker_notes": speaker_notes,
                "body_text": body_text,
                "inline_images": inline_images,
            }
        )
    return slides


# ---------------------------------------------------------------------------
# Rich-text helpers
# ---------------------------------------------------------------------------
def _add_formatted_runs(para, text, font_size=Pt(16), font_name=BODY_FONT,
                        color=DARK, bold=False, italic=False):
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            _add_italic_runs(para, part[2:-2], font_size, font_name, color, bold=True)
        else:
            _add_italic_runs(para, part, font_size, font_name, color, bold=bold)


def _add_italic_runs(para, text, font_size, font_name, color, bold=False):
    parts = re.split(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", text)
    for i, part in enumerate(parts):
        if not part:
            continue
        run = para.add_run()
        run.text = part
        run.font.size = font_size
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = bold
        if i % 2 == 1:
            run.font.italic = True


def _set_paragraph_text(para, text, font_size=Pt(16), font_name=BODY_FONT,
                        color=DARK, bold=False, alignment=None):
    para.clear()
    if alignment:
        para.alignment = alignment
    _add_formatted_runs(para, text, font_size, font_name, color, bold)


def _add_speaker_notes(slide, notes_text):
    if not notes_text:
        return
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.text = notes_text
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Slide building helpers
# ---------------------------------------------------------------------------
def _ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_title(slide, title_text, font_size=Pt(28), color=WHITE):
    ph = _ph(slide, 0)
    if ph is not None:
        _set_paragraph_text(ph.text_frame.paragraphs[0], title_text,
                            font_size=font_size, font_name=HEADING_FONT, color=color)
    else:
        _add_title_textbox(slide, title_text, font_size, color=color)


def _add_title_textbox(slide, title_text, font_size=Pt(28), color=WHITE):
    txBox = slide.shapes.add_textbox(Emu(588263), Emu(457200),
                                     Emu(11018520), Emu(553998))
    txBox.text_frame.word_wrap = True
    _set_paragraph_text(txBox.text_frame.paragraphs[0], title_text,
                        font_size=font_size, font_name=HEADING_FONT, color=color)


def _add_bullets_to_placeholder(slide, bullets, ph_idx=10):
    ph = _ph(slide, ph_idx)
    if ph is None:
        _add_bullets_textbox(slide, bullets)
        return
    tf = ph.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = 0
        para.space_after = Pt(4)
        _add_formatted_runs(para, bullet, font_size=Pt(16), font_name=BODY_FONT, color=WHITE)


def _add_bullets_textbox(slide, bullets, top=Emu(1435100)):
    txBox = slide.shapes.add_textbox(Emu(584200), top, Emu(11018838), Emu(4833938))
    txBox.text_frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = txBox.text_frame.paragraphs[0] if i == 0 else txBox.text_frame.add_paragraph()
        para.level = 0
        para.space_after = Pt(4)
        r = para.add_run()
        r.text = "• "
        r.font.size = Pt(16)
        r.font.name = BODY_FONT
        r.font.color.rgb = WHITE
        _add_formatted_runs(para, bullet, font_size=Pt(16), font_name=BODY_FONT, color=WHITE)


def _add_body_text(slide, body_lines, top=None):
    if not body_lines:
        return
    top = top or Emu(1435100)
    txBox = slide.shapes.add_textbox(Emu(584200), top, Emu(11018838), Emu(800000))
    txBox.text_frame.word_wrap = True
    for i, line in enumerate(body_lines):
        para = txBox.text_frame.paragraphs[0] if i == 0 else txBox.text_frame.add_paragraph()
        para.space_after = Pt(4)
        _add_formatted_runs(para, line, font_size=Pt(16), font_name=BODY_FONT, color=WHITE)


def _add_table(slide, table_data, top=Emu(1500000)):
    if not table_data:
        return
    rows = len(table_data)
    cols = len(table_data[0])
    left = Emu(584200)
    width = Emu(11018838)
    row_h = 370000
    height = Emu(row_h * rows)

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table

    # "No Style, No Grid" table design
    tbl_pr = tbl._tbl.tblPr
    tbl_pr.attrib["bandRow"] = "0"
    tbl_pr.attrib["bandCol"] = "0"
    tbl_pr.attrib["firstRow"] = "0"
    tbl_pr.attrib["lastRow"] = "0"
    tbl_pr.attrib["firstCol"] = "0"
    tbl_pr.attrib["lastCol"] = "0"
    style_id = tbl_pr.find(qn("a:tblStyleId"))
    if style_id is None:
        style_id = tbl_pr.makeelement(qn("a:tblStyleId"), {})
        tbl_pr.append(style_id)
    style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"

    col_w = width // cols
    for ci in range(cols):
        tbl.columns[ci].width = col_w

    for ri, row in enumerate(table_data):
        for ci, cell_text in enumerate(row):
            if ci >= cols:
                break
            cell = tbl.cell(ri, ci)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            is_hdr = ri == 0
            _add_formatted_runs(
                para, cell_text,
                font_size=Pt(12), font_name=BODY_FONT,
                color=WHITE if is_hdr else DARK, bold=is_hdr,
            )
            tcPr = cell._tc.get_or_add_tcPr()
            # Remove all cell borders for "No Style, No Grid"
            for border_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = tcPr.find(qn(border_tag))
                if ln is not None:
                    tcPr.remove(ln)
            # Clear any existing fills before setting new one
            for existing_fill in tcPr.findall(qn("a:solidFill")):
                tcPr.remove(existing_fill)
            sf = tcPr.makeelement(qn("a:solidFill"), {})
            clr = sf.makeelement(qn("a:srgbClr"), {"val": "0078D4" if is_hdr else "FFFFFF"})
            sf.append(clr)
            tcPr.append(sf)


def _add_code_textbox(slide, code_text, top=Emu(1500000)):
    margin = Emu(584200)
    width = Emu(11018838)
    line_count = code_text.count("\n") + 1
    height = max(Emu(500000), Emu(line_count * 200000))
    max_h = SLIDE_H - top - Emu(300000)
    if height > max_h:
        height = max_h

    txBox = slide.shapes.add_textbox(margin, top, width, height)
    txBox.text_frame.word_wrap = False
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = DARK  # Dark background for code blocks

    for i, line in enumerate(code_text.split("\n")):
        para = txBox.text_frame.paragraphs[0] if i == 0 else txBox.text_frame.add_paragraph()
        r = para.add_run()
        r.text = line
        r.font.size = Pt(10)
        r.font.name = CODE_FONT
        r.font.color.rgb = WHITE  # White text on dark background


# ---------------------------------------------------------------------------
# Main conversion
# ---------------------------------------------------------------------------
def convert(input_md: str, output_pptx: str, template_pptx: str,
            images_dir: str | None = None):
    print(f"Template : {template_pptx}")
    print(f"Input    : {input_md}")
    print(f"Output   : {output_pptx}")

    if images_dir is None:
        images_dir = os.path.join(os.path.dirname(input_md), "images")
    print(f"Images   : {images_dir}")

    prs = Presentation(template_pptx)

    # Delete all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    print(f"Cleared template ({len(prs.slides)} slides remain)")

    # Layout accessors — use branded layouts from reference deck
    def layout_content():
        return prs.slide_masters[MASTER_CONTENT].slide_layouts[LAYOUT_CONTENT_IDX]

    def layout_title():
        return prs.slide_masters[MASTER_CONTENT].slide_layouts[LAYOUT_TITLE_IDX]

    def layout_section():
        return prs.slide_masters[MASTER_SECTION].slide_layouts[LAYOUT_SECTION_IDX]

    def layout_divider():
        return prs.slide_masters[MASTER_DIVIDER].slide_layouts[LAYOUT_DIVIDER_IDX]

    def layout_notes():
        return prs.slide_masters[MASTER_NOTES].slide_layouts[LAYOUT_NOTES_IDX]

    def layout_closing():
        return prs.slide_masters[MASTER_NOTES].slide_layouts[LAYOUT_CLOSING_IDX]

    # Parse
    with open(input_md, "r", encoding="utf-8") as f:
        md_text = f.read()
    slides_data = parse_marp(md_text)
    print(f"Parsed {len(slides_data)} slides\n")

    for i, sd in enumerate(slides_data):
        stype = sd["type"]
        title = sd["title"]
        subtitle = sd["subtitle"]
        bullets = sd["bullets"]
        table = sd["table"]
        code_block = sd["code_block"]
        speaker_notes = sd["speaker_notes"]
        body_text = sd["body_text"]

        # --- choose layout ---
        if stype == "title":
            sl = layout_title()
            label = "Title Slide"
        elif stype == "lead":
            sl = layout_section()
            label = "Section Title"
        elif stype == "thankyou":
            sl = layout_closing()
            label = "Closing"
        elif stype == "notes":
            sl = layout_notes()
            label = "Notes"
        elif code_block:
            sl = layout_content()
            label = "Code"
        elif table and not bullets:
            sl = layout_content()
            label = "Table"
        elif table and bullets:
            sl = layout_content()
            label = "Table+Bullets"
        else:
            sl = layout_content()
            label = "Content"

        slide = prs.slides.add_slide(sl)
        short = (title[:40] + "..") if len(title) > 40 else title
        print(f"  {i+1:3d}  [{label:14s}]  {short}")

        # --- populate ---
        if stype == "title":
            _set_title(slide, title, Pt(32), color=BLUE)
            ph12 = _ph(slide, 12)
            if ph12 and subtitle:
                _set_paragraph_text(ph12.text_frame.paragraphs[0], subtitle,
                                    font_size=Pt(20), font_name=BODY_FONT, color=DARK)
            elif subtitle:
                tb = slide.shapes.add_textbox(Emu(584200), Emu(3810000),
                                              Emu(9144000), Emu(246221))
                _set_paragraph_text(tb.text_frame.paragraphs[0], subtitle,
                                    font_size=Pt(20), font_name=BODY_FONT, color=DARK)

        elif stype == "lead":
            # Render meaningful section title + subtitle — never just a section number
            _set_title(slide, title, Pt(36), color=WHITE)
            if subtitle:
                tb = slide.shapes.add_textbox(Emu(584200), Emu(2800000),
                                              Emu(11018838), Emu(600000))
                tb.text_frame.word_wrap = True
                _set_paragraph_text(tb.text_frame.paragraphs[0], subtitle,
                                    font_size=Pt(22), font_name=BODY_FONT, color=WHITE)

        elif stype == "thankyou":
            if title:
                _add_title_textbox(slide, title, Pt(36))

        elif stype == "notes":
            _set_title(slide, title, Pt(28))
            ph12 = _ph(slide, 12)
            if ph12 and bullets:
                tf = ph12.text_frame
                tf.clear()
                for bi, b in enumerate(bullets):
                    para = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                    para.space_after = Pt(4)
                    _add_formatted_runs(para, b, Pt(14), BODY_FONT, WHITE)
            elif bullets:
                _add_bullets_textbox(slide, bullets, top=Emu(1436688))

        elif code_block and table:
            _set_title(slide, title, Pt(28))
            _add_table(slide, table, top=Emu(1200000))
            tbl_h = Emu(370000 * len(table))
            _add_code_textbox(slide, code_block, top=Emu(1200000) + tbl_h + Emu(100000))

        elif code_block:
            _set_title(slide, title, Pt(28))
            ctop = Emu(1200000)
            if bullets:
                _add_bullets_textbox(slide, bullets, top=Emu(1100000))
                ctop = min(Emu(1100000 + len(bullets) * 280000 + 100000), Emu(3500000))
            if body_text:
                _add_body_text(slide, body_text, top=ctop)
                ctop = min(ctop + Emu(len(body_text) * 280000 + 100000), Emu(4000000))
            _add_code_textbox(slide, code_block, top=ctop)

        elif table:
            _set_title(slide, title, Pt(28))
            ttop = Emu(1200000)
            if bullets:
                _add_bullets_textbox(slide, bullets, top=Emu(1100000))
                ttop = min(Emu(1100000 + len(bullets) * 280000 + 200000), Emu(3000000))
            _add_table(slide, table, top=ttop)
            after = Emu(ttop + 370000 * len(table) + 100000)
            if body_text:
                _add_body_text(slide, body_text, top=after)

        else:
            _set_title(slide, title, Pt(28))
            if bullets:
                _add_bullets_to_placeholder(slide, bullets, ph_idx=10)
            if body_text:
                bt = Emu(1435100)
                if bullets:
                    bt = min(Emu(1435100 + len(bullets) * 280000 + 100000), Emu(5500000))
                _add_body_text(slide, body_text, top=bt)

        _add_speaker_notes(slide, speaker_notes)

        # --- embed image if mapped ---
        img_file = _image_for_slide(title)
        if img_file:
            img_path = os.path.join(images_dir, img_file)
            if os.path.isfile(img_path):
                if _is_diagram(img_file):
                    # Architecture diagram — replace code block area, centered
                    # Remove existing code block textbox if present
                    if code_block:
                        for shape in list(slide.shapes):
                            if shape.has_text_frame:
                                tf = shape.text_frame
                                if any(r.font.name == CODE_FONT
                                       for p in tf.paragraphs
                                       for r in p.runs if r.font.name):
                                    sp = shape._element
                                    sp.getparent().remove(sp)
                                    break
                    img_left = Inches(0.5)
                    img_top = Inches(2.0)
                    img_width = Inches(12.0)
                    max_img_h = SLIDE_H - img_top - Emu(200000)
                    with Image.open(img_path) as im:
                        iw, ih = im.size
                    aspect = ih / iw
                    img_height = int(img_width * aspect)
                    if img_height > max_img_h:
                        img_height = max_img_h
                        img_width = int(img_height / aspect)
                        img_left = (SLIDE_W - img_width) // 2
                    _add_image_to_slide(slide, img_path, img_left, img_top,
                                        img_width, img_height,
                                        alt_text=f"Architecture diagram: {title}")
                    print(f"        + diagram: {img_file}")
                else:
                    # Stock photo — add to the right side
                    img_width = Inches(4.0)
                    img_left = SLIDE_W - img_width - Inches(0.3)
                    img_top = Inches(2.0)
                    _add_image_to_slide(slide, img_path, img_left, img_top,
                                        width=img_width,
                                        alt_text=f"Illustration for: {title}")
                    print(f"        + photo: {img_file}")
            else:
                print(f"        ! missing: {img_file}")

        # --- fallback: embed inline markdown images if no IMAGE_MAP match ---
        if not img_file:
            inline_images = sd.get("inline_images", [])
            for alt_text, img_ref in inline_images:
                inline_path = os.path.join(images_dir, os.path.basename(img_ref))
                if not os.path.isfile(inline_path):
                    # Try relative to input markdown directory
                    inline_path = os.path.join(os.path.dirname(input_md), img_ref)
                if os.path.isfile(inline_path):
                    basename = os.path.basename(inline_path)
                    if _is_diagram(basename):
                        img_left = Inches(0.5)
                        img_top = Inches(2.4)
                        img_width = Inches(12.0)
                        max_img_h = SLIDE_H - img_top - Emu(200000)
                        with Image.open(inline_path) as im:
                            iw, ih = im.size
                        aspect = ih / iw
                        img_height = int(img_width * aspect)
                        if img_height > max_img_h:
                            img_height = max_img_h
                            img_width = int(img_height / aspect)
                        img_left = (SLIDE_W - img_width) // 2
                        _add_image_to_slide(slide, inline_path, img_left, img_top,
                                            img_width, img_height,
                                            alt_text=alt_text or f"Architecture diagram: {title}")
                        print(f"        + inline diagram: {basename}")
                    else:
                        img_width = Inches(4.0)
                        img_left = SLIDE_W - img_width - Inches(0.3)
                        img_top = Inches(2.0)
                        _add_image_to_slide(slide, inline_path, img_left, img_top,
                                            width=img_width,
                                            alt_text=alt_text or f"Illustration: {title}")
                        print(f"        + inline photo: {basename}")
                else:
                    print(f"        ! inline missing: {img_ref}")

    # Save
    os.makedirs(os.path.dirname(output_pptx) or ".", exist_ok=True)
    prs.save(output_pptx)
    size = os.path.getsize(output_pptx)
    print(f"\nSaved {len(prs.slides)} slides -> {output_pptx}  ({size:,} bytes)")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(
        description="Convert a Marp markdown deck to PowerPoint (.pptx)."
    )
    parser.add_argument("input", nargs="?", default=DEFAULT_INPUT,
                        help="Path to Marp .md file")
    parser.add_argument("output", nargs="?", default=DEFAULT_OUTPUT,
                        help="Path for output .pptx")
    parser.add_argument("--template", "-t", default=DEFAULT_TEMPLATE,
                        help="Path to reference PPTX template")
    parser.add_argument("--images-dir", default=None,
                        help="Directory with slide images (default: {input_dir}/images/)")
    args = parser.parse_args()

    convert(args.input, args.output, args.template, images_dir=args.images_dir)


if __name__ == "__main__":
    main()
